#!/usr/bin/env python
"""Sinh các fixture nhị phân cho validation corpus (tái tạo được, không check tay).

Vì sao có script này: các fixture text (`.txt/.md/.html`) đọc được trực tiếp nên cứ
commit thẳng; còn `.pptx/.xlsx/.png/.jpg/.pdf/.docx` là nhị phân — opaque khi review.
Script này dựng lại chúng từ nội dung khai báo ngay đây, nên fixture luôn khớp với
golden query trong manifest và ai cũng regenerate được:

    python eval/validation/generate_fixtures.py

Hai nhóm:
- Nhóm A (text-form, offline): pptx, xlsx — đi qua MarkItDown, KHÔNG cần OCR. Thêm vào
  `manifest.json` và suite offline hiện tại tự chạy.
- Nhóm B (OCR/vision, gated): png, jpg, PDF scan, docx-có-ảnh — chữ chỉ nằm trong ảnh,
  phải qua AI gateway. Khai trong `manifest_ocr.json`, test riêng chỉ chạy khi có provider.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

HERE = Path(__file__).resolve().parent


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def _font(size: int):
    """TrueType nếu có (chữ to, rõ cho OCR); fallback bitmap mặc định của PIL."""
    from PIL import ImageFont

    for candidate in (r"C:\Windows\Fonts\arial.ttf", "arial.ttf", "DejaVuSans.ttf"):
        try:
            return ImageFont.truetype(candidate, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _text_image(lines: list[str], *, size: int = 48, width: int = 1280, pad: int = 60):
    """Render nhiều dòng chữ đen trên nền trắng -> PIL.Image (để OCR đọc lại)."""
    from PIL import Image, ImageDraw

    font = _font(size)
    line_h = size + 22
    height = pad * 2 + line_h * len(lines)
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    for index, line in enumerate(lines):
        draw.text((pad, pad + index * line_h), line, fill="black", font=font)
    return image


# --------------------------------------------------------------------------- #
# Nhóm A — text-form (offline qua MarkItDown)
# --------------------------------------------------------------------------- #
def make_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Remote Work Policy"

    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tf = box.text_frame
    tf.text = "Eligible employees may work remotely up to three days per week."
    for line in (
        "Remote days must be agreed with the line manager in advance.",
        "Core collaboration hours remain 10:00 to 15:00 on remote days.",
        "Equipment and a secure VPN connection are required when working remotely.",
    ):
        tf.add_paragraph().text = line
    prs.save(str(path))


def make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Per Diem"
    rows = [
        ("Travel Per Diem Rates", "", ""),
        ("Category", "Allowance", "Currency"),
        ("Meals", "fifty per day", "USD"),
        ("Lodging", "one hundred eighty per night", "USD"),
        ("Ground transport", "thirty per day", "USD"),
        ("Note", "Daily meal allowance is fifty USD for domestic travel", ""),
    ]
    for row in rows:
        ws.append(row)
    wb.save(str(path))


# --------------------------------------------------------------------------- #
# Nhóm B — OCR/vision (gated)
# --------------------------------------------------------------------------- #
def make_png(path: Path) -> None:
    _text_image(
        [
            "GUEST WIFI ACCESS",
            "",
            "Network name: VFS-Guest",
            "Password: welcome2026",
            "Access expires after twenty four hours.",
        ]
    ).save(str(path), format="PNG")


def make_jpg(path: Path) -> None:
    _text_image(
        [
            "VISITOR PARKING NOTICE",
            "",
            "Visitor parking is on level B2 only.",
            "Register your plate at the front desk.",
            "Maximum stay is four hours.",
        ]
    ).save(str(path), format="JPEG", quality=92)


def make_scanned_pdf(path: Path) -> None:
    """PDF KHÔNG có text-layer: 1 trang là ảnh => ép parser rasterize -> OCR."""
    import fitz

    image = _text_image(
        [
            "FIRE EVACUATION PROCEDURE",
            "",
            "On hearing the alarm, leave the building immediately.",
            "Use the nearest staircase. Do not use the lifts.",
            "Assemble at the car park muster point.",
        ]
    )
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", quality=80)  # JPEG-stream giữ PDF nhỏ gọn
    doc = fitz.open()
    page = doc.new_page(width=image.width * 0.5, height=image.height * 0.5)
    page.insert_image(page.rect, stream=buffer.getvalue())
    doc.save(str(path), garbage=4, deflate=True)
    doc.close()


def make_docx_with_image(path: Path) -> None:
    """docx tối thiểu: 1 đoạn text + 1 ảnh nhúng trong word/media (text+image merge).

    Dựng zip thủ công vì python-docx không có trong env; `_read_docx_file` chỉ cần
    `word/document.xml` (text) và quét `word/media/*` (ảnh) — không soi relationships.
    """
    buffer = io.BytesIO()
    _text_image(
        [
            "EMERGENCY CONTACT CARD",
            "",
            "Facilities hotline: 1800 7000",
            "After hours security: extension 911",
        ]
    ).save(buffer, format="PNG")

    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="png" ContentType="image/png"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.'
        'openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        "<w:p><w:r><w:t>The full emergency contact list is shown in the card below.</w:t>"
        "</w:r></w:p>"
        "<w:p><w:r><w:t>Keep this card visible at every workstation.</w:t></w:r></w:p>"
        "</w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("word/document.xml", document_xml)
        archive.writestr("word/media/image1.png", buffer.getvalue())


# --------------------------------------------------------------------------- #
def main() -> None:
    builders = {
        "remote_work_policy.pptx": make_pptx,
        "travel_per_diem.xlsx": make_xlsx,
        "guest_wifi.png": make_png,
        "visitor_parking.jpg": make_jpg,
        "fire_evacuation_scanned.pdf": make_scanned_pdf,
        "emergency_contacts.docx": make_docx_with_image,
    }
    for name, builder in builders.items():
        target = HERE / name
        builder(target)
        print(f"wrote {target.relative_to(HERE.parents[1])} ({target.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
